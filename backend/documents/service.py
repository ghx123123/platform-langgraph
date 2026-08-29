import io
import base64
import hashlib
import os
import re
import statistics
import unicodedata
import zipfile
from collections import Counter
from dataclasses import dataclass, field
from functools import lru_cache
from math import ceil
from pathlib import Path
from typing import Any, Callable
from uuid import uuid4

from backend.workflows.models import DocumentSection, KnowledgePoint


ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".md", ".txt"}
MAX_DOCUMENT_BYTES = 30 * 1024 * 1024
MAX_DOCUMENT_CHARACTERS = 120000
MAX_DOCUMENT_SECTIONS = 160

# S4: MinerU 解析引擎 — 本机已有部署, 复用不下载. 可用 环境变量 覆盖.
MINERU_BIN = Path(os.environ.get("TC_MINERU_BIN") or "D:/software/anaconda/envs/mineru/Scripts/mineru.exe")
MINERU_WORK_DIR = Path(os.environ.get("TC_MINERU_WORK") or "")

# 扫描版教材常见 OCR 误识别（按短语精确匹配，避免误伤正常文字）
_OCR_FIXES: tuple[tuple[str, str], ...] = (
    ("Pyhon", "Python"), ("卩hon", "Python"), ("Pyhom", "Python"), ("Pyt hon", "Python"),
    ("逆辑", "逻辑"), ("干分位", "千分位"), ("下画线", "下划线"), ("二雄码", "二维码"),
    ("字符申", "字符串"), ("返曰", "返回"), ("题回", "返回"),
    ("mapO", "map()"), ("reduceO", "reduce()"), ("rangeO", "range()"), ("ilterO", "filter()"),
    ("2.Python 常用内置对象", "2.1 Python 常用内置对象"), ("2.13字符串", "2.1.3 字符串"),
    ("基本输人输出", "基本输入输出"), ("选代", "迭代"), ("range()的数", "range()函数"),
    ("的数、", "、"), ("内数", "函数"),
)

# 章节辅助标题，不作为知识点
_SKIP_HEADINGS = ("学习目标", "本章小结", "本章习题", "小结", "习题", "思考与练习", "本章介绍", "本章重点", "本章难点", "补充说明", "注意事项")


@dataclass
class ExtractedDocument:
    text: str
    engine: str
    page_count: int = 0
    title_count: int = 0
    table_count: int = 0
    image_count: int = 0
    text_block_count: int = 0
    scanned_page_count: int = 0
    ocr_page_count: int = 0
    ocr_image_count: int = 0
    source_line_count: int = 0
    noisy_line_count: int = 0
    pages: list[str] = field(default_factory=list)
    page_reports: list[dict[str, Any]] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


def _normalize(text: str) -> str:
    text = unicodedata.normalize("NFKC", text)
    text = text.replace("\x00", "").replace("\u00ad", "").replace("\u200b", "")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _markdown_table(rows: list[list[Any]]) -> str:
    cleaned = [[re.sub(r"\s+", " ", str(cell or "")).strip().replace("|", "\\|") for cell in row] for row in rows]
    cleaned = [row for row in cleaned if any(row)]
    if not cleaned:
        return ""
    width = max(len(row) for row in cleaned)
    cleaned = [row + [""] * (width - len(row)) for row in cleaned]
    header = cleaned[0]
    return "\n".join([
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
        *("| " + " | ".join(row) + " |" for row in cleaned[1:]),
    ])


def _quality_report(extension: str, extracted: ExtractedDocument, text: str) -> dict[str, Any]:
    score = 100
    warnings = list(dict.fromkeys(extracted.warnings))
    if len(text) < 500:
        score -= 18
        warnings.append("提取文本较少，请确认材料是否主要由图片构成")
    replacement_ratio = text.count("�") / max(len(text), 1)
    if replacement_ratio > 0.001:
        score -= 20
        warnings.append("检测到较多无法解码字符")
    if extracted.page_count and extracted.scanned_page_count and not extracted.ocr_page_count:
        ratio = extracted.scanned_page_count / extracted.page_count
        score -= min(45, round(ratio * 45))
        warnings.append(f"{extracted.scanned_page_count} 页未提取到足够文字，可能需要 OCR")
    elif extracted.ocr_page_count:
        # OCR restores searchable text, but not native-text fidelity for
        # formulas, diagrams, or low-confidence characters.
        ocr_ratio = extracted.ocr_page_count / max(extracted.page_count, 1)
        score -= 8 + round(ocr_ratio * 8)
        warnings.append(f"原始材料包含 {extracted.scanned_page_count} 页扫描图像，已用中文 OCR 识别 {extracted.ocr_page_count} 页；公式、图表和低清文字建议人工复核")
    if extracted.title_count == 0:
        score -= 8
        warnings.append("未识别到明确标题层级，目录将按页或自然段生成")
    if extension == ".pptx" and extracted.page_count:
        average_characters = len(text) / extracted.page_count
        if extracted.image_count >= extracted.page_count and average_characters < 200:
            score -= 16
            warnings.append("课件以图片和图示为主，当前主要分析可访问文字；图中公式、关系和标注需要人工复核")
    if extension == ".docx" and extracted.image_count >= 8:
        characters_per_image = len(text) / extracted.image_count
        if extracted.ocr_image_count:
            score -= 16
            warnings.append(
                f"已识别 {extracted.ocr_image_count} 张嵌入图片的文字；复杂图示、连线关系和公式仍需要人工复核"
            )
        elif characters_per_image < 400:
            score -= 16
            warnings.append(
                f"文档包含 {extracted.image_count} 张图片且可访问文字相对较少；截图、图中标注和视觉关系需要人工复核"
            )
    noise_ratio = extracted.noisy_line_count / max(extracted.source_line_count, 1)
    if noise_ratio >= 0.03:
        score -= min(32, round(noise_ratio * 120))
        warnings.append(f"约 {round(noise_ratio * 100)}% 的文本行疑似包含 OCR 或编码噪声")
    if extracted.page_count > 1 and extracted.text_block_count < extracted.page_count:
        score -= 10
    score = max(0, min(100, score))
    level = "high" if score >= 85 else "medium" if score >= 65 else "low"
    return {
        "format": extension.lstrip(".").upper(),
        "engine": extracted.engine,
        "quality_score": score,
        "quality_level": level,
        "page_count": extracted.page_count,
        "title_count": extracted.title_count,
        "table_count": extracted.table_count,
        "image_count": extracted.image_count,
        "text_block_count": extracted.text_block_count,
        "scanned_page_count": extracted.scanned_page_count,
        "ocr_page_count": extracted.ocr_page_count,
        "ocr_image_count": extracted.ocr_image_count,
        "page_reports": extracted.page_reports,
        "warnings": list(dict.fromkeys(warnings))[:20],
    }


def fix_ocr_text(text: str) -> str:
    text = text.replace("\u63d0\u4e66", "\u63d0\u51fa").replace("\u5377\u4ece", "\u5377\u79ef").replace("\u5173\u5065", "\u5173\u952e")
    """修正扫描版教材常见的 OCR 误识别，降低下游分析与生成噪声。"""
    for old, new in _OCR_FIXES:
        text = text.replace(old, new)
    return text


def _docx_text(data: bytes) -> ExtractedDocument:
    """按 Word 正文中的段落/表格顺序提取，并用样式恢复标题层级。"""
    from docx import Document
    from docx.table import Table

    document = Document(io.BytesIO(data))
    blocks: list[str] = []
    title_count = 0
    table_count = 0
    text_block_count = 0
    seen_text: set[str] = set()

    for item in document.iter_inner_content():
        if isinstance(item, Table):
            table = _markdown_table([[cell.text for cell in row.cells] for row in item.rows])
            if table:
                blocks.append(table)
                table_count += 1
                text_block_count += 1
            continue
        value = re.sub(r"\s+", " ", item.text).strip()
        if not value:
            continue
        style_name = (item.style.name if item.style else "") or ""
        heading = re.search(r"(?:Heading|标题)\s*([1-6])", style_name, re.IGNORECASE)
        if heading:
            value = f"{'#' * int(heading.group(1))} {value}"
            title_count += 1
        elif "Title" in style_name or style_name == "标题":
            value = f"# {value}"
            title_count += 1
        elif "List" in style_name or "列表" in style_name:
            value = f"- {value}"
        blocks.append(value)
        seen_text.add(item.text.strip())
        text_block_count += 1

    # 文本框不属于 document.paragraphs，单独补充，避免复杂课件式 Word 丢字。
    text_boxes: list[str] = []
    for box in document._element.xpath(".//w:txbxContent"):
        for paragraph in box.xpath(".//w:p"):
            value = "".join(paragraph.xpath(".//w:t/text()"))
            value = re.sub(r"\s+", " ", value).strip()
            if value and value not in seen_text:
                text_boxes.append(value)
                seen_text.add(value)
    if text_boxes:
        blocks.extend(["## 文本框内容", *text_boxes])
        title_count += 1
        text_block_count += len(text_boxes)

    page_count = 0
    warnings: list[str] = []
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        if "docProps/app.xml" in archive.namelist():
            match = re.search(rb"<Pages>(\d+)</Pages>", archive.read("docProps/app.xml"))
            page_count = int(match.group(1)) if match else 0
        for part_name, label in (("word/footnotes.xml", "脚注"), ("word/endnotes.xml", "尾注")):
            if part_name not in archive.namelist():
                continue
            part = archive.read(part_name).decode("utf-8", errors="ignore")
            entries = [re.sub(r"<[^>]+>", "", entry).strip() for entry in re.findall(r"<w:p[\s\S]*?</w:p>", part)]
            entries = [re.sub(r"\s+", " ", entry) for entry in entries if entry]
            if entries:
                blocks.extend([f"## {label}", *entries])
                title_count += 1
                text_block_count += len(entries)
        if any(name.startswith("word/charts/") for name in archive.namelist()):
            warnings.append("文档包含图表，已提取相邻说明文字，但未还原图形视觉关系")

    if page_count == 0:
        warnings.append("文档未保存可靠页数元数据，页数显示为未知；正文、表格和标题仍按文档顺序提取")

    native_text = "\n\n".join(blocks)
    enhanced_text, enhanced_tables, ocr_image_count, enhancement_warnings, enhancement_ocr_engine = _markitdown_docx_text(data)
    use_enhanced = bool(enhanced_text and len(enhanced_text) > len(native_text) * 1.02)
    final_text = enhanced_text if use_enhanced else native_text
    warnings.extend(enhancement_warnings)

    return ExtractedDocument(
        text=final_text,
        engine=(
            f"python-docx + MarkItDown + {enhancement_ocr_engine or 'OCR'}"
            if use_enhanced else "python-docx"
        ),
        page_count=page_count,
        title_count=title_count,
        table_count=max(table_count, enhanced_tables if use_enhanced else 0),
        image_count=len(document.inline_shapes),
        text_block_count=text_block_count + (ocr_image_count if use_enhanced else 0),
        ocr_image_count=ocr_image_count if use_enhanced else 0,
        source_line_count=len(final_text.splitlines()),
        noisy_line_count=sum(int(_looks_noisy_line(line)) for line in final_text.splitlines() if line.strip()),
        warnings=warnings,
    )


def _span_line_text(spans: list[dict[str, Any]]) -> str:
    value = ""
    for span in spans:
        part = str(span.get("text", "")).strip()
        if not part:
            continue
        if value and value[-1].isascii() and value[-1].isalnum() and part[0].isascii() and part[0].isalnum():
            value += " "
        value += part
    return value.strip()


def _looks_noisy_line(value: str) -> bool:
    compact = re.sub(r"\s+", "", value)
    if len(compact) < 3:
        return False
    chinese = len(re.findall(r"[\u4e00-\u9fff]", compact))
    letters = len(re.findall(r"[A-Za-z]", compact))
    digits = len(re.findall(r"\d", compact))
    if digits / len(compact) > 0.55 and len(compact) > 12:
        return True
    mixed_tokens = re.findall(r"\b(?=\w*[A-Za-z])(?=\w*\d)\w{6,}\b", value)
    if mixed_tokens:
        return True
    if chinese == 0 and letters >= 8:
        words = re.findall(r"[A-Za-z]{4,}", value)
        if words and all(word.islower() for word in words) and not re.search(r"\b(python|boolean|integer|float|string|list|tuple|dict|set|none|true|false)\b", value, re.IGNORECASE):
            return True
    if chinese and chinese / len(compact) < 0.18 and letters + digits > chinese * 3:
        return True
    return False


def _looks_like_structural_heading(value: str) -> bool:
    value = value.strip(" ：:")
    if not 2 <= len(value) <= 80 or _looks_noisy_line(value):
        return False
    if re.fullmatch(r"第\s*\d+\s*页", value):
        return False
    # 编程/技术标题里 [as 别名] 等方括号是合法语法，不拒绝; 仍拒绝 = <> {} 与 >> 等比较符特征
    if re.search(r"[=<>{}]", value) or ">>" in value:
        return False
    if re.search(r"[。；，,;！？!?]$", value):
        return False
    action_prefixes = (
        "使用", "支持", "内置函数", "把", "如果", "返回", "测试", "连接", "创建", "比较",
        "注意", "等价", "查看", "不支持", "指定", "转换", "生成", "提取", "修改", "按", "先",
        "用于", "一个", "这里", "后面", "直接", "自动", "不影响", "不对", "访问", "组合", "代码注释",
    )
    return not value.startswith(action_prefixes)


def _looks_like_markdown_heading(value: str) -> bool:
    """宽松版 markdown 标题判定：信任 `#` 提供的结构化标题信息。

    MinerU/转换器输出的 markdown 标题已是作者的结构化标记，不再用
    「_looks_noisy_line 的代码特征/英文比例」二次过滤(那会误杀
    1.1Python语言简介 等编程教材标题)。只做显式噪声剔除。
    """
    value = value.strip(" ：:")
    if not 2 <= len(value) <= 80:
        return False
    if re.fullmatch(r"第\s*\d+\s*页", value):
        return False
    if re.fullmatch(r"(?:END|THE END|目录|目次|TOC|地址|网址)\s*", value, re.IGNORECASE | re.X):
        return False
    # 明显非标题：句末标点/纯符号行
    if re.search(r"[。；！？!?]\s*$", value):
        return False
    if not re.search(r"[一-鿿A-Za-z0-9]+", value):
        return False
    return True


def _looks_like_pdf_heading(value: str) -> bool:
    value = value.strip()
    if not _looks_like_structural_heading(value):
        return False
    if re.search(r"(用来|可以|支持|返回|表示|例如|其中|需要|如果|由于|进行|具有)", value) and len(value) > 12:
        return False
    chinese = len(re.findall(r"[\u4e00-\u9fff]", value))
    if chinese:
        return chinese >= 3 and chinese / max(len(re.sub(r"\s+", "", value)), 1) >= 0.3
    words = re.findall(r"[A-Za-z]+", value)
    return len(words) >= 2 and any(word[:1].isupper() for word in words)


def _looks_like_ocr_heading(value: str, line_index: int = 0) -> bool:
    """Use conservative heuristics for slide OCR because scanned pages have no font metadata."""
    value = value.strip().lstrip("●○口O•·- ")
    if not 3 <= len(value) <= 60 or _looks_noisy_line(value):
        return False
    if re.match(r"^第\s*\d+\s*[章节]", value):
        return True
    numbered = re.match(r"^(\d+(?:\.\d+){1,3})\s*(.+)$", value)
    if numbered and _is_section_number(numbered.group(1)) and _looks_like_numbered_title(numbered.group(2)):
        return True
    if re.search(r"[。；，、：？！.!?,;:]", value) or re.search(r"[=+*/%<>{}\[\]]", value):
        return False
    chinese = len(re.findall(r"[\u4e00-\u9fff]", value))
    if line_index == 0 and (chinese >= 3 or len(re.findall(r"[A-Za-z]+", value)) >= 2):
        return True
    return False


@lru_cache(maxsize=1)
def _cnocr_engine() -> Any:
    """Load the OCR model once per worker so repeated uploads stay responsive."""
    from cnocr import CnOcr

    return CnOcr()


@lru_cache(maxsize=1)
def _rapidocr_engine() -> Any:
    """Load PP-OCRv6 once; it is faster and more accurate on Chinese textbooks."""
    from rapidocr import RapidOCR

    return RapidOCR()


def _preferred_ocr_engine() -> tuple[str, Any, list[str]]:
    """Prefer RapidOCR while retaining CnOCR for installations without ONNX models."""
    try:
        return "RapidOCR v6", _rapidocr_engine(), []
    except Exception as rapid_exc:
        try:
            return "CnOCR", _cnocr_engine(), [
                f"RapidOCR v6 未能加载，已回退 CnOCR：{type(rapid_exc).__name__}"
            ]
        except Exception as cnocr_exc:
            raise RuntimeError(
                f"RapidOCR v6 ({type(rapid_exc).__name__}) 与 CnOCR "
                f"({type(cnocr_exc).__name__}) 均未能加载"
            ) from cnocr_exc


def _normalized_ocr_results(engine: Any, image: Any) -> list[tuple[Any, str, float]]:
    """Normalize RapidOCR and CnOCR output to position/text/confidence tuples."""
    if callable(engine):
        output = engine(image)
    else:
        output = engine.ocr(image)

    if hasattr(output, "txts"):
        texts = list(output.txts) if output.txts is not None else []
        scores = list(output.scores) if output.scores is not None else []
        boxes = list(output.boxes) if output.boxes is not None else []
        return [
            (
                boxes[index] if index < len(boxes) else None,
                str(value),
                float(scores[index]) if index < len(scores) else 1.0,
            )
            for index, value in enumerate(texts)
        ]

    normalized: list[tuple[Any, str, float]] = []
    for result in output or []:
        if not isinstance(result, dict):
            continue
        normalized.append((
            result.get("position"),
            str(result.get("text", "")),
            float(result.get("score", 1.0) or 0.0),
        ))
    return normalized


def _full_page_image(page: Any) -> tuple[int, float] | None:
    """Return the dominant image xref when an image covers nearly the whole page."""
    page_area = max(float(page.rect.get_area()), 1.0)
    dominant: tuple[int, float] | None = None
    for image in page.get_images(full=True):
        xref = int(image[0])
        try:
            coverage = max((float(rect.get_area()) / page_area for rect in page.get_image_rects(xref)), default=0.0)
        except (RuntimeError, ValueError):
            coverage = 0.0
        if dominant is None or coverage > dominant[1]:
            dominant = (xref, coverage)
    return dominant if dominant and dominant[1] >= 0.90 else None


def _has_hidden_ocr_overlay(page: Any) -> bool:
    """Detect scanned pages whose existing OCR text is invisible and often unreliable."""
    if _full_page_image(page) is None:
        return False
    try:
        traces = page.get_texttrace()
    except (RuntimeError, ValueError):
        return False
    text_traces = [trace for trace in traces if trace.get("chars")]
    return bool(text_traces) and sum(trace.get("type") == 3 for trace in text_traces) / len(text_traces) >= 0.90


def _ocr_page_elements(page: Any, engine: Any) -> list[dict[str, Any]]:
    """OCR a page while preserving boxes, confidence and visual reading order."""
    import numpy as np
    from PIL import Image

    dominant = _full_page_image(page)
    image = None
    if dominant is not None:
        try:
            source = page.parent.extract_image(dominant[0]).get("image", b"")
            if source:
                image = np.asarray(Image.open(io.BytesIO(source)).convert("RGB"))
        except (OSError, RuntimeError, ValueError):
            image = None
    if image is None:
        pixmap = page.get_pixmap(matrix=__import__("fitz").Matrix(2.0, 2.0), alpha=False)
        image = np.frombuffer(pixmap.samples, dtype=np.uint8).reshape(pixmap.height, pixmap.width, pixmap.n)

    results = _normalized_ocr_results(engine, image)
    image_height, image_width = image.shape[:2]
    x_scale = float(page.rect.width) / max(image_width, 1)
    y_scale = float(page.rect.height) / max(image_height, 1)
    positioned: list[dict[str, Any]] = []
    for index, (position, text, score) in enumerate(results):
        value = re.sub(r"\s+", " ", text).strip()
        if not value or score < 0.08:
            continue
        try:
            points = position.tolist() if hasattr(position, "tolist") else position
            xs = [float(point[0]) for point in points]
            ys = [float(point[1]) for point in points]
            x0, x1 = min(xs) * x_scale, max(xs) * x_scale
            y0, y1 = min(ys) * y_scale, max(ys) * y_scale
        except (TypeError, ValueError, IndexError):
            x0, x1 = 0.0, float(page.rect.width)
            y0, y1 = float(index * 12), float(index * 12 + 10)
        positioned.append({
            "y": y0,
            "x": x0,
            "text": value,
            "size": max(1.0, y1 - y0),
            "bbox": (x0, y0, x1, y1),
            "score": score,
            "ocr": True,
        })

    row_height = max(1.5, y_scale * 6)
    positioned.sort(key=lambda item: (round(item["y"] / row_height), item["x"]))
    for line_index, item in enumerate(positioned):
        item["line_index"] = line_index
    return positioned


def _ocr_page_lines(page: Any, engine: Any) -> list[str]:
    """Compatibility wrapper returning visually ordered OCR text lines."""
    return [item["text"] for item in _ocr_page_elements(page, engine)]


def _ocr_image_bytes(data: bytes, engine: Any) -> list[str]:
    """OCR one embedded Office image while bounding memory and inference cost."""
    import numpy as np
    from PIL import Image

    image = Image.open(io.BytesIO(data)).convert("RGB")
    width, height = image.size
    if width * height < 80000 or min(width, height) < 120:
        return []
    if max(width, height) > 2200:
        scale = 2200 / max(width, height)
        image = image.resize((round(width * scale), round(height * scale)), Image.Resampling.LANCZOS)
    results = _normalized_ocr_results(engine, np.asarray(image))
    positioned: list[tuple[float, float, str]] = []
    for index, (position, text, score) in enumerate(results):
        value = re.sub(r"\s+", " ", text).strip()
        if not value or score < 0.12:
            continue
        try:
            points = position.tolist() if hasattr(position, "tolist") else position
            y = min(float(point[1]) for point in points)
            x = min(float(point[0]) for point in points)
        except (TypeError, ValueError, IndexError):
            y, x = float(index), 0.0
        positioned.append((y, x, value))
    return [value for _, _, value in sorted(positioned, key=lambda item: (round(item[0] / 6), item[1]))]


def _docx_image_blobs_in_order(data: bytes) -> list[bytes]:
    """Read embedded image bytes in the same drawing order as the DOCX body."""
    from docx import Document
    from docx.oxml.ns import qn

    document = Document(io.BytesIO(data))
    images: list[bytes] = []
    for element in document.element.body.iter():
        if element.tag != qn("a:blip"):
            continue
        relationship_id = element.get(qn("r:embed"))
        if not relationship_id:
            continue
        related_part = document.part.related_parts.get(relationship_id)
        blob = getattr(related_part, "blob", None)
        if blob:
            images.append(blob)
    return images


def _markitdown_docx_text(data: bytes) -> tuple[str, int, int, list[str], str]:
    """Use MarkItDown for reading order and place OCR text where images occur."""
    try:
        import warnings

        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", message="Couldn't find ffmpeg or avconv.*", category=RuntimeWarning)
            from markitdown import MarkItDown
    except ImportError:
        return "", 0, 0, ["MarkItDown 未安装，已使用内置 Word 解析器"]

    try:
        markdown = MarkItDown(enable_plugins=False).convert_stream(
            io.BytesIO(data), file_extension=".docx"
        ).text_content
    except Exception as exc:
        return "", 0, 0, [f"MarkItDown 增强解析失败，已回退内置解析器：{type(exc).__name__}"], ""

    # MarkItDown 0.1.x intentionally renders DOCX images as
    # ``data:image/...;base64...`` placeholders. Recover those bytes from the
    # OpenXML relationships while retaining MarkItDown's reading order.
    image_pattern = re.compile(
        r"!\[([^\]]*)\]\(data:(image/[^;()\s]+);base64(?:(?:,([A-Za-z0-9+/=\r\n]+))|\.\.\.)\)"
    )
    matches = list(image_pattern.finditer(markdown))
    if not matches:
        cleaned = re.sub(r"\n{3,}", "\n\n", markdown).strip()
        tables = sum(bool(re.match(r"^\|\s*:?-{3,}", line)) for line in cleaned.splitlines())
        return cleaned, tables, 0, ["已使用开源 MarkItDown 补充 Word 标题、表格与阅读顺序"], ""

    try:
        ocr_engine_name, engine, engine_warnings = _preferred_ocr_engine()
    except Exception as exc:
        engine = None
        ocr_engine_name = ""
        engine_warnings = []
        ocr_warning = f"嵌入图片 OCR 未启用：{type(exc).__name__}"
    else:
        ocr_warning = ""

    ordered_images = _docx_image_blobs_in_order(data)
    image_position = 0
    cache: dict[bytes, list[str]] = {}
    recognized_hashes: set[bytes] = set()
    attempted_hashes: set[bytes] = set()
    max_unique_images = 40

    def replace_image(match: re.Match[str]) -> str:
        nonlocal image_position
        alt = re.sub(r"\s+", " ", match.group(1)).strip() or "文档图片"
        encoded = match.group(3)
        if encoded:
            try:
                image_data = base64.b64decode(encoded, validate=False)
            except (ValueError, TypeError):
                image_data = b""
        else:
            image_data = ordered_images[image_position] if image_position < len(ordered_images) else b""
        image_position += 1
        if not image_data:
            return f"\n\n[图片：{alt}]\n\n"
        digest = hashlib.sha256(image_data).digest()
        if digest not in cache:
            if engine is None or len(attempted_hashes) >= max_unique_images:
                cache[digest] = []
            else:
                attempted_hashes.add(digest)
                try:
                    cache[digest] = _ocr_image_bytes(image_data, engine)
                except Exception:
                    cache[digest] = []
        lines = cache[digest]
        if lines:
            recognized_hashes.add(digest)
            return f"\n\n图片识别（{alt}）：\n" + "\n".join(lines) + "\n\n"
        return f"\n\n[图片：{alt}]\n\n"

    enhanced = image_pattern.sub(replace_image, markdown)
    enhanced = re.sub(r"<!--.*?-->", "", enhanced, flags=re.DOTALL)
    enhanced = re.sub(r"\n{3,}", "\n\n", enhanced).strip()
    table_count = sum(bool(re.match(r"^\|\s*:?-{3,}", line)) for line in enhanced.splitlines())
    warnings = [
        f"已使用开源 MarkItDown 补充 Word 标题、表格和图片位置，并用 {ocr_engine_name or 'OCR'} "
        f"识别 {len(recognized_hashes)} 张嵌入图片"
    ]
    warnings.extend(engine_warnings)
    if ocr_warning:
        warnings.append(ocr_warning)
    if len(matches) != len(ordered_images):
        warnings.append(
            f"文档正文包含 {len(matches)} 个图片位置，OpenXML 读取到 {len(ordered_images)} 张图片；未匹配图片需结合原页复核"
        )
    if len(attempted_hashes) >= max_unique_images and len(cache) > max_unique_images:
        warnings.append(f"嵌入图片超过 {max_unique_images} 张，仅对前 {max_unique_images} 张执行 OCR")
    return enhanced, table_count, len(recognized_hashes), warnings, ocr_engine_name


def _mineru_pdf_text(data: bytes, progress_cb: Callable[[int, str], None] | None = None) -> ExtractedDocument:
    """用 MinerU (本机已部署 magic-pdf) 版面识别提取整个 PDF 为 Markdown 文本。

    复用 D:/software/anaconda/envs/mineru 的 mineru.exe, 不下载模型; 输出 markdown 读回。
    引擎较慢(版面识别), 仅当用户显式选择 engine='mineru' 时使用。
    progress_cb(percent, message) 用于把 MinerU 子进程的 tqdm 进度回传给调用方。
    """
    import subprocess
    import tempfile

    if not MINERU_BIN.exists():
        raise RuntimeError(f"MinerU 未找到: {MINERU_BIN}")
    with tempfile.TemporaryDirectory() as tmp:
        pdf_path = Path(tmp) / "input.pdf"
        pdf_path.write_bytes(data)
        out_dir = Path(tmp) / "out"
        out_dir.mkdir(exist_ok=True)
        proc = subprocess.Popen(
            [str(MINERU_BIN), "-p", str(pdf_path), "-o", str(out_dir)],
            stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True,
            encoding="utf-8", errors="replace", bufsize=1,
        )
        # 进度解析: MinerU(tqdm) 输出形如 "Layout Predict:  67%|██████...  6/9 ..."
        # 每行用 \r 或行结束刷新; 解析百分比并映射阶段到 0-100。
        if progress_cb is not None:
            progress_cb(5, "正在启动 MinerU 版面识别…")
        last_pct = 5
        stage_weights = {"init": (0, 5), "Layout": (5, 35), "MFD": (35, 60), "MFR": (60, 80), "Generate": (80, 100)}
        tail = ""
        assert proc.stdout is not None
        while True:
            chunk = proc.stdout.read(4096)
            if not chunk:
                break
            tail += chunk
            tail = tail.replace("\r", "\n")
            segments = tail.split("\n")
            tail = segments[-1]
            for segment in segments[-6:]:
                match = re.search(r"(Layout|MFD|MFR|Generate|Format)[^%]*?(\d{1,3})\s*%", segment)
                if not match:
                    continue
                stage = match.group(1)
                pct = int(match.group(2))
                lo, hi = stage_weights.get(stage, (5, 100))
                overall = lo + int((hi - lo) * pct / 100)
                if overall > last_pct:
                    last_pct = overall
                    if progress_cb is not None:
                        progress_cb(max(1, min(99, overall)), f"MinerU {stage} 识别 {pct}%")
        proc.wait()
        if proc.returncode != 0:
            if progress_cb is not None:
                progress_cb(99, "MinerU 失败，正在清理")
            raise RuntimeError(f"MinerU 失败: {(proc.stderr or proc.stdout or '') if proc.stderr else ''}")
        md_files = sorted(out_dir.rglob("*.md"))
        if not md_files:
            raise RuntimeError("MinerU 未产出 Markdown 结果")
        text = md_files[0].read_text(encoding="utf-8", errors="replace").strip()
        if len(text) < 10:
            raise RuntimeError("MinerU 提取文本过短")
        if progress_cb is not None:
            progress_cb(100, "MinerU 识别完成")
        return ExtractedDocument(
            text=text,
            engine="MinerU (magic-pdf layout)",
            page_count=len(list(out_dir.rglob("*.pdf"))) or 1,
            text_block_count=1,
        )


def _pdf_text(data: bytes) -> ExtractedDocument:
    """利用 PDF 坐标、字号和表格边界恢复阅读顺序，过滤重复页眉页脚。"""
    import fitz

    document = fitz.open(stream=data, filetype="pdf")
    page_records: list[dict[str, Any]] = []
    margin_counter: Counter[str] = Counter()
    table_detection = len(document) <= 120
    table_count = 0
    image_count = 0
    scanned_pages = 0
    hidden_overlay_pages = 0
    text_block_count = 0
    source_line_count = 0
    noisy_line_count = 0

    for page in document:
        height = page.rect.height
        hidden_ocr = _has_hidden_ocr_overlay(page)
        hidden_overlay_pages += int(hidden_ocr)
        dictionary = page.get_text("dict", sort=True)
        elements: list[dict[str, Any]] = []
        page_characters = 0
        sizes: list[float] = []
        table_boxes: list[tuple[float, float, float, float]] = []
        tables: list[Any] = []
        if table_detection and not hidden_ocr:
            try:
                tables = list(page.find_tables().tables)
                table_boxes = [tuple(table.bbox) for table in tables]
                table_count += len(tables)
            except (ValueError, RuntimeError):
                tables = []

        for block in dictionary.get("blocks", []):
            if block.get("type") != 0 or hidden_ocr:
                continue
            text_block_count += 1
            for line in block.get("lines", []):
                spans = line.get("spans", [])
                value = _span_line_text(spans)
                if not value:
                    continue
                source_line_count += 1
                noisy_line_count += int(_looks_noisy_line(value))
                bbox = tuple(line.get("bbox", (0, 0, 0, 0)))
                center = ((bbox[0] + bbox[2]) / 2, (bbox[1] + bbox[3]) / 2)
                if any(box[0] <= center[0] <= box[2] and box[1] <= center[1] <= box[3] for box in table_boxes):
                    continue
                size = max((float(span.get("size", 0)) for span in spans), default=0)
                sizes.extend(float(span.get("size", 0)) for span in spans if span.get("text", "").strip())
                page_characters += len(value)
                elements.append({"y": bbox[1], "x": bbox[0], "text": value, "size": size, "bbox": bbox})
                if bbox[1] < height * 0.09 or bbox[3] > height * 0.92:
                    key = re.sub(r"\d+", "#", re.sub(r"\s+", "", value))
                    if 2 <= len(key) <= 80:
                        margin_counter[key] += 1

        for table in tables:
            markdown = _markdown_table(table.extract())
            if markdown:
                elements.append({"y": table.bbox[1], "x": table.bbox[0], "text": markdown, "size": 0, "bbox": table.bbox, "table": True})
        if page_characters < 30 or hidden_ocr:
            scanned_pages += 1
        image_count += len(page.get_images(full=True))
        page_records.append({
            "height": height,
            "elements": elements,
            "sizes": sizes,
            "scanned": page_characters < 30 or hidden_ocr,
            "source_kind": "hidden_ocr" if hidden_ocr else "scanned" if page_characters < 30 else "native",
            "ocr_confidence": None,
        })

    ocr_page_count = 0
    ocr_engine_name = ""
    ocr_warnings: list[str] = []
    if scanned_pages:
        try:
            ocr_engine_name, ocr_engine, fallback_warnings = _preferred_ocr_engine()
            ocr_warnings.extend(fallback_warnings)
        except Exception as exc:  # pragma: no cover - depends on optional local OCR install
            ocr_engine = None
            ocr_warnings.append(
                f"扫描页未启用 OCR：{type(exc).__name__}。请安装 rapidocr、onnxruntime 或 cnocr 后重试"
            )
        if ocr_engine is not None:
            for index, page in enumerate(document):
                if not page_records[index]["scanned"]:
                    continue
                try:
                    ocr_elements = _ocr_page_elements(page, ocr_engine)
                except Exception as exc:  # pragma: no cover - depends on page/image data
                    ocr_warnings.append(f"第 {index + 1} 页 OCR 失败：{type(exc).__name__}")
                    continue
                if not ocr_elements:
                    continue
                ocr_lines = [element["text"] for element in ocr_elements]
                page_records[index]["elements"] = ocr_elements
                page_records[index]["sizes"] = [element["size"] for element in ocr_elements]
                page_records[index]["ocr_confidence"] = statistics.mean(
                    element["score"] for element in ocr_elements
                )
                text_block_count += 1
                source_line_count += len(ocr_lines)
                noisy_line_count += sum(int(_looks_noisy_line(line)) for line in ocr_lines)
                ocr_page_count += 1

    repeated_threshold = max(2, ceil(len(document) * 0.35))
    repeated_margins = {key for key, count in margin_counter.items() if count >= repeated_threshold}
    pages: list[str] = []
    page_reports: list[dict[str, Any]] = []
    title_count = 0
    for page_index, page_record in enumerate(page_records):
        body_size = statistics.median(page_record["sizes"]) if page_record["sizes"] else 10
        lines: list[str] = []
        page_title_count = 0
        previous = ""
        sorted_elements = sorted(page_record["elements"], key=lambda item: (round(item["y"] / 3), item["x"]))
        pending_number = ""
        for element in sorted_elements:
            value = element["text"].strip()
            key = re.sub(r"\d+", "#", re.sub(r"\s+", "", value))
            bbox = element["bbox"]
            in_margin = bbox[1] < page_record["height"] * 0.09 or bbox[3] > page_record["height"] * 0.92
            if in_margin and key in repeated_margins:
                continue
            if re.fullmatch(r"[-—–·•\s]*\d{1,4}[-—–·•\s]*", value):
                continue
            if value == previous:
                continue
            previous = value
            if element.get("table"):
                lines.append(value)
                continue
            if value.startswith("#"):
                comment = value.lstrip("#").strip()
                if comment:
                    lines.append(f"代码注释：{comment}")
                continue
            if _is_section_number(value):
                pending_number = value
                continue
            if pending_number:
                if _looks_like_pdf_heading(value):
                    value = f"## {pending_number} {value}"
                    title_count += 1
                    pending_number = ""
                    lines.append(value)
                    continue
                lines.append(pending_number)
                pending_number = ""
            ocr_heading = bool(
                element.get("ocr")
                and _looks_like_ocr_heading(
                    value,
                    element.get("line_index", 99) if page_index == 0 else 99,
                )
            )
            is_heading = ocr_heading or (
                element["size"] >= body_size * 1.28 and _looks_like_pdf_heading(value)
            )
            if is_heading:
                level = 1 if element["size"] >= body_size * 1.65 else 2
                value = f"{'#' * level} {value}"
                title_count += 1
                page_title_count += 1
            lines.append(value)
        if pending_number:
            lines.append(pending_number)
        page_text = "\n".join(lines).strip()
        pages.append(page_text)
        page_reports.append({
            "page_number": page_index + 1,
            "character_count": len(page_text),
            "line_count": len([line for line in lines if line.strip()]),
            "title_count": page_title_count,
            "ocr_applied": page_record["ocr_confidence"] is not None,
            "ocr_confidence": (
                round(float(page_record["ocr_confidence"]), 4)
                if page_record["ocr_confidence"] is not None else None
            ),
            "source_kind": page_record["source_kind"],
        })

    if title_count == 0 and len(pages) > 1:
        pages = [f"## 第 {index} 页\n{page}" for index, page in enumerate(pages, start=1) if page]
    warnings = [*ocr_warnings]
    if hidden_overlay_pages:
        warnings.append(
            f"检测到 {hidden_overlay_pages} 页整页扫描图像带隐藏 OCR 文字层，已忽略旧文字层并逐页重新识别"
        )
    if not table_detection:
        warnings.append("PDF 超过 120 页，已优先保证全文阅读顺序，未逐页执行表格边界识别")
    if repeated_margins:
        warnings.append(f"已过滤 {len(repeated_margins)} 类重复页眉或页脚")
    return ExtractedDocument(
        text="\n\n".join(page for page in pages if page),
        engine=f"PyMuPDF + {ocr_engine_name}" if ocr_page_count else "PyMuPDF",
        page_count=len(document),
        title_count=title_count,
        table_count=table_count,
        image_count=image_count,
        text_block_count=text_block_count,
        scanned_page_count=scanned_pages,
        ocr_page_count=ocr_page_count,
        source_line_count=source_line_count,
        noisy_line_count=noisy_line_count,
        pages=pages,
        page_reports=page_reports,
        warnings=warnings,
    )


def _pptx_text(data: bytes) -> ExtractedDocument:
    """按幻灯片和版面位置提取文本、表格、图表与讲者备注。"""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(io.BytesIO(data))
    slides: list[str] = []
    table_count = 0
    image_count = 0
    text_block_count = 0
    chart_count = 0
    note_count = 0
    actual_title_count = 0
    thin_slide_count = 0

    def chart_markdown(shape: Any) -> str:
        nonlocal chart_count
        chart = shape.chart
        categories = [str(item.label) for item in chart.plots[0].categories] if chart.plots and chart.plots[0].categories else []
        rows: list[list[Any]] = [["系列", *categories]]
        for series in chart.series:
            rows.append([series.name, *list(series.values)])
        chart_count += 1
        return "### 图表数据\n" + _markdown_table(rows) if len(rows) > 1 else ""

    def shape_lines(shape: Any) -> list[str]:
        nonlocal table_count, image_count, text_block_count
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            nested = sorted(getattr(shape, "shapes", []), key=lambda item: (item.top, item.left))
            return [line for item in nested for line in shape_lines(item)]
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1
        if getattr(shape, "has_chart", False):
            value = chart_markdown(shape)
            return [value] if value else []
        if getattr(shape, "has_table", False):
            table_count += 1
            value = _markdown_table([[cell.text for cell in row.cells] for row in shape.table.rows])
            return [value] if value else []
        if getattr(shape, "has_text_frame", False):
            result: list[str] = []
            for paragraph in shape.text_frame.paragraphs:
                value = "".join(run.text for run in paragraph.runs).strip() or paragraph.text.strip()
                if not value:
                    continue
                prefix = "  " * min(paragraph.level, 4) + ("- " if paragraph.level or len(shape.text_frame.paragraphs) > 1 else "")
                result.append(prefix + value)
                text_block_count += 1
            return result
        return []

    for index, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        actual_title = title_shape.text.strip() if title_shape and title_shape.text.strip() else ""
        body: list[str] = []
        for shape in sorted(slide.shapes, key=lambda item: (round(item.top / 10000), item.left)):
            if shape is title_shape:
                continue
            body.extend(shape_lines(shape))
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except (AttributeError, ValueError):
            notes = ""
        if notes:
            note_lines = [line.strip() for line in notes.splitlines() if line.strip() and line.strip() != title]
            if note_lines:
                body.extend(["### 讲者备注", *note_lines])
                note_count += 1
        unique_body = list(dict.fromkeys(line for line in body if line.strip() and line.strip("- ") != actual_title))
        if not actual_title:
            for body_index, line in enumerate(unique_body):
                candidate = line.lstrip("- ").strip()
                if _looks_like_structural_heading(candidate):
                    actual_title = candidate
                    unique_body.pop(body_index)
                    break
        if actual_title:
            actual_title_count += 1
        title = actual_title or "无标题"
        slide_heading = title if title_shape else f"幻灯片 {index} | {title}"
        if sum(len(line) for line in unique_body) < 30:
            thin_slide_count += 1
        slides.append("\n".join([f"## {slide_heading}", *unique_body]))

    warnings = []
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        if any(name.startswith("ppt/diagrams/") for name in archive.namelist()):
            warnings.append("课件包含 SmartArt，已提取可访问文字，但图形关系建议人工复核")
    if chart_count:
        warnings.append(f"已将 {chart_count} 个图表转换为可分析的数据表")
    if note_count:
        warnings.append(f"已纳入 {note_count} 页讲者备注")
    if thin_slide_count:
        warnings.append(f"{thin_slide_count} 页提取到的可访问文字少于 30 字，图片中的文字或图形关系可能未被识别")
    return ExtractedDocument(
        text="\n\n".join(slides),
        engine="python-pptx",
        page_count=len(presentation.slides),
        title_count=actual_title_count,
        table_count=table_count + chart_count,
        image_count=image_count,
        text_block_count=text_block_count,
        warnings=warnings,
    )


def _course_name(text: str, filename: str) -> str:
    lines = text.splitlines()
    for index, line in enumerate(lines):
        stripped = line.strip().lstrip("# ").strip()
        chapter = re.search(r"第\s*[一二三四五六七八九十0-9]+\s*[章节][^|#\n]{0,60}", stripped)
        if chapter:
            value = chapter.group(0).strip()
            if re.fullmatch(r"第\s*[一二三四五六七八九十0-9]+\s*[章节]", value):
                for following in lines[index + 1:index + 3]:
                    candidate = following.strip().lstrip("# ").strip(" ：:")
                    if candidate and not re.search(r"^(教材|作者|出版社)\s*[：:]?", candidate):
                        value = f"{value} {candidate}"
                        break
            if 4 <= len(value) <= 80:
                return value
    """优先取文档标题（Markdown 或“第X章”式教材标题），回退到文件名。"""
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("#"):
            title = stripped.lstrip("#").strip(" ：:")
            if 2 <= len(title) <= 60:
                return title
        if re.match(r"^第[一二三四五六七八九十0-9]+章", stripped) and len(stripped) <= 60:
            return stripped
        break
    return Path(filename).stem


def _is_section_number(value: str) -> bool:
    if not re.fullmatch(r"\d+(?:\.\d+){1,3}", value):
        return False
    parts = [int(part) for part in value.split(".")]
    return parts[0] <= 50 and all(part <= 99 for part in parts[1:])


def _looks_like_numbered_title(value: str) -> bool:
    value = value.strip(" ：:")
    if not 2 <= len(value) <= 80 or re.search(r"[=<>{}]", value) or ">>" in value:
        return False
    if re.search(r"[。；！？!?]$", value):
        return False
    return len(re.findall(r"[\u4e00-\u9fff]", value)) >= 2


def _is_numbered_heading(line: str) -> bool:
    match = re.match(r"^([0-9]+(?:\.[0-9]+){1,3})(?:\s+|(?=[\u4e00-\u9fff]))(.+)$", line)
    return bool(match and _is_section_number(match.group(1)) and _looks_like_numbered_title(match.group(2)))


def _is_heading(line: str) -> bool:
    s = line.strip()
    # Markdown 标题要求：# + 空格 + 标题，且不含代码特征（括号/等号/句末标点），避免把 Python 代码注释当标题
    markdown = re.match(r"^#{1,6}\s+(.+)$", s)
    if markdown and _looks_like_markdown_heading(markdown.group(1)):
        return True
    if re.match(r"^第[一二三四五六七八九十0-9]+章", s):
        return True
    if _is_numbered_heading(s):
        return True
    if any(word in s for word in ("重点", "难点", "知识点", "核心")):
        return True
    return False


def _is_code_noise(line: str) -> bool:
    """识别 Python 代码行与 OCR 噪声行：解释器提示符、纯 ASCII、符号堆砌等。"""
    s = line.strip()
    if not s:
        return True
    if ">>" in s:  # Python 交互式解释器提示符（OCR 常把 >>> 识别成 ?>> 或 > >）
        return True
    if not re.search(r"[\u4e00-\u9fff]", s):
        return True
    chinese = len(re.findall(r"[\u4e00-\u9fff]", s))
    if chinese / max(len(s), 1) < 0.25 and not _is_numbered_heading(s):
        return True
    if re.search(r"[=+*/%<>{}]", s) and "，" not in s and not _is_numbered_heading(s):
        return True
    if re.search(r"[■□◆●◉★☆▲△]", s):
        return True
    return False


def extract_knowledge_points(text: str) -> list[KnowledgePoint]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    selected: list[str] = []
    for raw in lines:
        if _is_code_noise(raw):
            continue
        if not _is_heading(raw):
            continue
        title = re.sub(r"^[#\s\-•*]+", "", raw).strip(" ：:")
        # 剥离“第X章”/“N.N(.N)”编号前缀，保留知识点名称
        title = re.sub(r"^第[一二三四五六七八九十0-9]+章[：:\s]*", "", title)
        title = title.lstrip("●○口O•·- ").strip()
        title = re.sub(r"^[0-9]+(\.[0-9]+){1,2}\s*", "", title)
        # 清除行尾 OCR 残留符号
        title = re.sub(r"[\s回囗□]+$", "", title).strip(" ：:")
        if not title or any(skip in title for skip in _SKIP_HEADINGS) or re.fullmatch(r"第\s*\d+\s*页", title):
            continue
        if re.search(r"(?:Introduction|\u6559\u6750|\u51fa\u7248\u793e|\u738b\u4e07\u826f|\u5de5\u4e1a\u4eba\u5de5\u667a\u80fd\u5bfc\u8bba|\u5de5\u4e1a\u76d1\u8054\u7f51|\d{4}\u5e74|\u5206\u949f|END|\u63a5\u53d7\u5230|\u89c6\u9891)", title, re.IGNORECASE):
            continue
        if 4 <= len(title) <= 40:
            selected.append(title)
    if not selected:
        selected = [
            re.sub(r"^[\-•*\d.、]+", "", line).strip(" ：:")
            for line in lines
            if 5 <= len(line) <= 40 and not _is_code_noise(line)
        ]
    unique: list[str] = []
    seen_keys: set[str] = set()
    for item in selected:
        key = re.sub(r"\s+", "", item).lower()
        if item and key not in seen_keys:
            seen_keys.add(key)
            unique.append(item)
    course_title = _course_name(text, "")
    unique = [item for item in unique if item != course_title][:12]
    key_quota = max(1, round(len(unique) / 3))
    return [
        KnowledgePoint(
            title=title,
            chapter="课程材料",
            is_key_point=any(word in title for word in ("重点", "核心", "关键")) or index < key_quota,
            difficulty_level="较难" if any(word in title for word in ("难", "综合", "推导")) else "中等",
        )
        for index, title in enumerate(unique)
    ]


def _slide_number_ranges(numbers: list[int]) -> str:
    ranges: list[str] = []
    start = previous = numbers[0]
    for number in numbers[1:]:
        if number == previous + 1:
            previous = number
            continue
        ranges.append(str(start) if start == previous else f"{start}-{previous}")
        start = previous = number
    ranges.append(str(start) if start == previous else f"{start}-{previous}")
    return "、".join(ranges)


def extract_pptx_knowledge_points(text: str) -> list[KnowledgePoint]:
    """Collapse repeated slide headings into teacher-selectable topic ranges."""
    grouped: dict[str, tuple[str, list[int]]] = {}
    slide_position = 0
    for raw in text.splitlines():
        heading = re.match(r"^##\s+(.+)$", raw.strip())
        if not heading:
            continue
        slide_position += 1
        value = heading.group(1).strip()
        explicit = re.match(r"^幻灯片\s+(\d+)\s*\|\s*(.+)$", value)
        slide_number = int(explicit.group(1)) if explicit else slide_position
        title = explicit.group(2).strip() if explicit else value
        title = re.sub(r"\s+解答多样化$", "", title).strip(" ：:")
        if not title or any(skip in title for skip in _SKIP_HEADINGS):
            continue
        if re.fullmatch(r"第\s*[一二三四五六七八九十百0-9]+\s*章.*", title):
            continue
        if re.fullmatch(r"第\s*\d+\s*节", title) or re.search(r"视频|END|人工智能导论$", title, re.IGNORECASE):
            continue
        key = re.sub(r"[\s：:]", "", title).lower()
        if key not in grouped:
            grouped[key] = (title, [])
        grouped[key][1].append(slide_number)

    topics = list(grouped.values())[:12]
    if not topics:
        return extract_knowledge_points(text)
    key_quota = max(1, round(len(topics) / 3))
    return [
        KnowledgePoint(
            title=f"{title}（幻灯片 {_slide_number_ranges(numbers)}）",
            chapter="课程材料",
            is_key_point=index < key_quota,
            difficulty_level="中等",
        )
        for index, (title, numbers) in enumerate(topics)
    ]


def _section_heading(line: str) -> tuple[str, int] | None:
    stripped = line.strip()
    markdown = re.match(r"^(#{1,6})\s+(.+)$", stripped)
    if markdown and _looks_like_markdown_heading(markdown.group(2)):
        return markdown.group(2).strip(" ：:"), len(markdown.group(1))
    if re.match(r"^第[一二三四五六七八九十百0-9]+[章节篇单元]", stripped):
        return stripped.strip(" ：:"), 1
    numbered = re.match(r"^([0-9]+(?:\.[0-9]+){1,3})(?:\s+|(?=[\u4e00-\u9fff]))([^=<>]{2,80})$", stripped)
    if numbered and _is_section_number(numbered.group(1)) and _looks_like_numbered_title(numbered.group(2)):
        level = min(6, numbered.group(1).count(".") + 2)
        return stripped.strip(" ：:"), level
    return None


def _section_preview(content: str, title: str, limit: int = 180) -> str:
    lines = [line.strip() for line in content.splitlines() if line.strip()]
    if lines and lines[0].lstrip("#").strip(" ：:") == title.lstrip("#").strip(" ：:"):
        lines = lines[1:]
    compact = re.sub(r"\s+", " ", " ".join(lines)).strip()
    return compact[:limit] + ("…" if len(compact) > limit else "")


def extract_document_sections(
    text: str,
    preserve_duplicate_titles: bool = False,
    page_offsets: list[tuple[int, int]] | None = None,
) -> list[DocumentSection]:
    """从标题层级构建原文索引；无标题文档按自然段分块，保证全文可预览。"""
    markers: list[tuple[int, str, int]] = []
    offset = 0
    for line in text.splitlines(keepends=True):
        if preserve_duplicate_titles:
            slide_heading = re.match(r"^##\s+(.+)$", line.strip())
            heading = (slide_heading.group(1).strip(" ：:"), 2) if slide_heading else None
        else:
            heading = _section_heading(line)
        if heading:
            title = heading[0].lstrip("●○口O•·- ").strip()
            outline_noise = r"(?:Introduction|\u6559\u6750|\u51fa\u7248\u793e|\u738b\u4e07\u826f|\u5de5\u4e1a\u4eba\u5de5\u667a\u80fd\u5bfc\u8bba|\u5de5\u4e1a\u76d1\u8054\u7f51|\d{4}\u5e74|\u5206\u949f|END|\u63a5\u53d7\u5230|\u89c6\u9891|TTE)"
            if preserve_duplicate_titles or (
                not re.search(outline_noise, title, re.IGNORECASE)
            ):
                markers.append((offset, title, heading[1]))
        offset += len(line)

    ranges: list[tuple[int, int, str, int]] = []
    if markers:
        if not preserve_duplicate_titles:
            unique_markers: list[tuple[int, str, int]] = []
            seen_titles: set[str] = set()
            for marker in markers:
                key = re.sub(r"\s+", "", marker[1]).lower()
                if key in seen_titles:
                    continue
                seen_titles.add(key)
                unique_markers.append(marker)
            markers = unique_markers
        markers = markers[:MAX_DOCUMENT_SECTIONS]
        first_offset = markers[0][0]
        if first_offset > 0 and text[:first_offset].strip():
            ranges.append((0, first_offset, "文档导言", 1))
        for index, (start, title, level) in enumerate(markers):
            end = markers[index + 1][0] if index + 1 < len(markers) else len(text)
            ranges.append((start, end, title, level))
    else:
        chunk_size = 2400
        start = 0
        while start < len(text) and len(ranges) < MAX_DOCUMENT_SECTIONS:
            target = min(len(text), start + chunk_size)
            end = text.find("\n\n", target)
            if end == -1 or end - target > 500:
                end = target
            else:
                end += 2
            ranges.append((start, end, f"内容片段 {len(ranges) + 1}", 1))
            start = end

    sections: list[DocumentSection] = []
    for index, (start, end, title, level) in enumerate(ranges, start=1):
        content = text[start:end].strip()
        if not content:
            continue
        page_start = page_end = None
        if page_offsets:
            start_candidates = [number for number, (page_start_offset, _) in enumerate(page_offsets, start=1) if page_start_offset <= start]
            end_probe = max(start, end - 1)
            end_candidates = [number for number, (page_start_offset, _) in enumerate(page_offsets, start=1) if page_start_offset <= end_probe]
            page_start = start_candidates[-1] if start_candidates else 1
            page_end = end_candidates[-1] if end_candidates else page_start
        sections.append(DocumentSection(
            id=f"section-{index}",
            title=title[:120],
            level=level,
            start_offset=start,
            end_offset=end,
            character_count=len(content),
            preview=_section_preview(content, title),
            page_start=page_start,
            page_end=page_end,
        ))
    return sections


def parse_document(filename: str, data: bytes, engine: str = "mineru", progress_cb: Callable[[int, str], None] | None = None) -> dict:
    extension = Path(filename).suffix.lower()
    if extension not in ALLOWED_EXTENSIONS:
        raise ValueError("仅支持 PDF、DOCX、PPTX、Markdown 和 TXT 文件")
    if not data:
        raise ValueError("上传文件为空")
    if len(data) > MAX_DOCUMENT_BYTES:
        raise ValueError("文档不能超过 30 MB")
    if extension == ".docx":
        extracted = _docx_text(data)
    elif extension == ".pptx":
        extracted = _pptx_text(data)
    elif extension == ".pdf":
        # S5: 默认 MinerU 版面识别; RapidOCR 仅当显式选择 rapidocr 时
        if engine == "mineru":
            extracted = _mineru_pdf_text(data, progress_cb=progress_cb)
        else:
            extracted = _pdf_text(data)
    else:
        extracted = ExtractedDocument(
            text=data.decode("utf-8-sig", errors="replace"),
            engine="UTF-8 text",
            page_count=1,
            text_block_count=1,
        )
    page_offsets: list[tuple[int, int]] | None = None
    if extracted.pages:
        normalized_pages = [fix_ocr_text(_normalize(page)) for page in extracted.pages]
        chunks: list[str] = []
        page_offsets = []
        cursor = 0
        for page_index, page in enumerate(normalized_pages):
            if chunks:
                chunks.append("\n\n")
                cursor += 2
            start = cursor
            chunks.append(page)
            cursor += len(page)
            page_offsets.append((start, cursor))
            if page_index < len(extracted.page_reports):
                extracted.page_reports[page_index]["character_count"] = len(page)
                extracted.page_reports[page_index]["line_count"] = len([line for line in page.splitlines() if line.strip()])
        text = "".join(chunks).strip()
    else:
        text = fix_ocr_text(_normalize(extracted.text))
    if len(text) < 10:
        raise ValueError("未能从文档中提取足够的文本内容")
    stored_text = text[:MAX_DOCUMENT_CHARACTERS]
    points = extract_pptx_knowledge_points(stored_text) if extension == ".pptx" else extract_knowledge_points(stored_text)
    if extension == ".pptx":
        source_chapter = re.search(r"第\s*(\d+)\s*章", Path(filename).stem)
        content_chapter = re.search(r"第\s*(\d+)\s*章", text)
        if source_chapter and content_chapter and source_chapter.group(1) != content_chapter.group(1):
            extracted.warnings.append(
                f"文件名标注第 {source_chapter.group(1)} 章，但课件正文首先识别为第 {content_chapter.group(1)} 章，请确认材料版本"
            )
    stored_page_offsets = None
    if page_offsets:
        stored_page_offsets = [
            (min(start, len(stored_text)), min(end, len(stored_text)))
            for start, end in page_offsets
            if start < len(stored_text)
        ]
    sections = extract_document_sections(
        stored_text,
        preserve_duplicate_titles=extension == ".pptx",
        page_offsets=stored_page_offsets,
    )
    return {
        "document_id": str(uuid4()),
        "file_name": Path(filename).name,
        "course_name": _course_name(text, filename),
        "raw_text": stored_text,
        "knowledge_points": [point.model_dump() for point in points],
        "sections": [section.model_dump() for section in sections],
        "extraction_report": _quality_report(extension, extracted, text),
        "character_count": len(text),
        "processed_character_count": len(stored_text),
        "is_truncated": len(text) > len(stored_text),
        # 页偏移: 供前端"页对齐"视图 (raw_text 按页切分)
        "page_offsets": stored_page_offsets,
    }
