"""文档解析器回归测试：OCR 纠错、课程名识别、知识点抽取、代码/噪声行过滤。"""
import io

import fitz
from docx import Document
from pptx import Presentation
from backend.documents.service import (
    _course_name,
    extract_document_sections,
    extract_knowledge_points,
    fix_ocr_text,
    parse_document,
)

SAMPLE = """第2章 内置对象、运算符、表达式、关键字
Pyhon内置对象不需要安装和导入任何模块就可以直接使用。
本章学习目标
掌握运算符的用法
2.2 Python 运算符与表达式
2.2.3 成员测试运算符
>>>x = [1, 2, 3]
>>>print(x)
2.2.5 逆辑运算符
逆辑运算符用来连接条件表达式
2.3 Python 常用内置函数用法
2.3.2 最值与求和
# 这里要连续按(Enter)键两次才能执行
2.3.6 mapO的数、reduceO的数、ilterO函数
2.4 Python 关键字简要说明
本章小结
本章习题
"""


def test_fix_ocr_text_common_errors():
    fixed = fix_ocr_text("Pyhon和卩hon都是误识别，字符申要修正")
    assert "Python" in fixed
    assert "字符串" in fixed


def test_course_name_recognizes_chapter_heading():
    assert _course_name(SAMPLE, "book.pdf") == "第2章 内置对象、运算符、表达式、关键字"


def test_course_name_falls_back_to_filename():
    assert _course_name("普通正文，没有标题", "my_course.pdf") == "my_course"


def test_course_name_combines_split_chapter_number_and_title():
    assert _course_name("第10章\n深度学习与大语言模型\n教材：工业人工智能导论", "book.pdf") == "第10章 深度学习与大语言模型"


def test_knowledge_points_filter_code_and_ocr_noise():
    points = extract_knowledge_points(fix_ocr_text(SAMPLE))
    titles = [p.title for p in points]
    assert "Python 运算符与表达式" in titles
    assert "成员测试运算符" in titles
    assert "逻辑运算符" in titles  # OCR：逆辑 -> 逻辑
    assert "最值与求和" in titles
    assert "map()、reduce()、filter()函数" in titles
    assert "Python 关键字简要说明" in titles
    # 代码行与注释行不得入选
    assert not any("x = [1, 2, 3]" in t for t in titles)
    assert not any("print(x)" in t for t in titles)
    assert not any("Enter" in t for t in titles)
    assert not any("本章" in t for t in titles)


def test_parse_document_txt_end_to_end():
    result = parse_document("course.txt", SAMPLE.encode("utf-8"))
    assert result["course_name"] == "第2章 内置对象、运算符、表达式、关键字"
    assert result["character_count"] > 0
    assert result["knowledge_points"]
    assert "Pyhon" not in result["raw_text"]
    assert "Python" in result["raw_text"]
    assert result["sections"]
    assert result["processed_character_count"] == result["character_count"]
    assert result["is_truncated"] is False


def test_document_sections_map_back_to_original_text():
    text = "# 第一部分\n概念说明。\n\n## 第二部分\n应用示例。"
    sections = extract_document_sections(text)

    assert [section.title for section in sections] == ["第一部分", "第二部分"]
    assert text[sections[1].start_offset:sections[1].end_offset].startswith("## 第二部分")
    assert sections[0].preview == "概念说明。"


def test_parse_document_pptx_preserves_slide_outline_and_body():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "热力学第一定律"
    slide.placeholders[1].text = "内能变化与热量、做功的关系\n符号约定与适用条件"
    second = presentation.slides.add_slide(presentation.slide_layouts[1])
    second.shapes.title.text = "典型过程分析"
    second.placeholders[1].text = "等容过程\n绝热过程"
    stream = io.BytesIO()
    presentation.save(stream)

    result = parse_document("热力学课件.pptx", stream.getvalue())

    assert result["course_name"] == "热力学第一定律"
    assert "## 热力学第一定律" in result["raw_text"]
    assert "符号约定与适用条件" in result["raw_text"]
    assert "## 典型过程分析" in result["raw_text"]
    assert result["knowledge_points"]
    assert [section["title"] for section in result["sections"]] == ["热力学第一定律", "典型过程分析"]
    assert result["extraction_report"]["engine"] == "python-pptx"
    assert result["extraction_report"]["page_count"] == 2
    assert result["extraction_report"]["title_count"] == 2


def test_parse_document_pptx_preserves_repeated_slide_titles():
    presentation = Presentation()
    for body in ("第一组案例", "第二组案例", "第三组案例"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "监督学习"
        slide.placeholders[1].text = body
    stream = io.BytesIO()
    presentation.save(stream)

    result = parse_document("第7章机器学习.pptx", stream.getvalue())

    assert len(result["sections"]) == 3
    assert [section["title"] for section in result["sections"]] == ["监督学习"] * 3


def test_parse_document_pptx_uses_text_box_title_for_blank_layout():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(100, 100, 4000000, 500000)
    box.text = "机器学习的一般步骤"
    body = slide.shapes.add_textbox(100, 800000, 4000000, 1500000)
    body.text = "数据准备、模型训练和模型评估构成完整流程。"
    stream = io.BytesIO()
    presentation.save(stream)

    result = parse_document("机器学习.pptx", stream.getvalue())

    assert result["extraction_report"]["title_count"] == 1
    assert len(result["sections"]) == 1
    assert result["sections"][0]["title"].startswith("幻灯片 1 | 机器学习的一般步骤")


def test_parse_document_pptx_uses_slide_boundaries_not_body_section_labels():
    presentation = Presentation()
    for title, body_text in (("课程目录", "第1节\n第2节"), ("监督学习", "7.1 定义\n分类任务")):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body_text
    stream = io.BytesIO()
    presentation.save(stream)

    result = parse_document("第7章机器学习.pptx", stream.getvalue())

    assert len(result["sections"]) == 2
    assert [section["title"] for section in result["sections"]] == ["课程目录", "监督学习"]


def test_parse_document_pptx_does_not_drop_video_slide_title():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "教学视频"
    slide.placeholders[1].text = "机器学习应用案例"
    stream = io.BytesIO()
    presentation.save(stream)

    result = parse_document("机器学习.pptx", stream.getvalue())

    assert len(result["sections"]) == 1
    assert result["sections"][0]["title"] == "教学视频"


def test_parse_document_pptx_collapses_repeated_topic_titles_into_slide_ranges():
    presentation = Presentation()
    for title in ["机器学习的发展", "机器学习的发展", "机器学习的发展", "监督学习"]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(0, 0, 4_000_000, 800_000)
        textbox.text = title
    stream = io.BytesIO()
    presentation.save(stream)

    result = parse_document("机器学习.pptx", stream.getvalue())
    titles = [point["title"] for point in result["knowledge_points"]]

    assert titles == ["机器学习的发展（幻灯片 1-3）", "监督学习（幻灯片 4）"]
    assert len(result["sections"]) == 4


def test_parse_document_docx_preserves_heading_and_table_order():
    document = Document()
    document.add_heading("热力学基础", level=1)
    document.add_paragraph("能量守恒是分析热力过程的主线。")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "过程"
    table.cell(0, 1).text = "特征"
    table.cell(1, 0).text = "等容"
    table.cell(1, 1).text = "体积不变"
    document.add_heading("符号约定", level=2)
    stream = io.BytesIO()
    document.save(stream)

    result = parse_document("热力学.docx", stream.getvalue())

    assert result["raw_text"].index("# 热力学基础") < result["raw_text"].index("| 过程 | 特征 |")
    assert result["raw_text"].index("| 过程 | 特征 |") < result["raw_text"].index("## 符号约定")
    assert result["extraction_report"]["table_count"] == 1
    assert result["extraction_report"]["title_count"] == 2
    assert [section["title"] for section in result["sections"]] == ["热力学基础", "符号约定"]


def test_image_heavy_docx_requires_review_and_does_not_invent_page_count():
    document = Document()
    document.add_heading("智能备课支撑材料", level=1)
    document.add_paragraph("平台流程说明" * 80)
    stream = io.BytesIO()
    document.save(stream)

    from backend.documents.service import ExtractedDocument, _quality_report

    extracted = ExtractedDocument(
        text="平台流程说明" * 100,
        engine="python-docx",
        page_count=0,
        title_count=1,
        image_count=12,
        text_block_count=2,
    )
    report = _quality_report(".docx", extracted, extracted.text)

    assert report["page_count"] == 0
    assert report["quality_level"] == "medium"
    assert report["quality_score"] == 84
    assert any("图片" in warning and "人工复核" in warning for warning in report["warnings"])


def test_docx_uses_longer_markitdown_text_and_reports_image_ocr(monkeypatch):
    document = Document()
    document.add_heading("课程说明", level=1)
    document.add_paragraph("内置解析正文")
    stream = io.BytesIO()
    document.save(stream)
    enhanced = "# 课程说明\n\n" + "开源解析补充正文" * 40 + "\n\n图片识别：流程图文字"

    monkeypatch.setattr(
        "backend.documents.service._markitdown_docx_text",
        lambda data: (enhanced, 3, 2, ["已完成增强解析"], "RapidOCR v6"),
    )

    result = parse_document("课程说明.docx", stream.getvalue())

    assert "开源解析补充正文" in result["raw_text"]
    assert result["extraction_report"]["engine"] == "python-docx + MarkItDown + RapidOCR v6"
    assert result["extraction_report"]["table_count"] == 3
    assert result["extraction_report"]["ocr_image_count"] == 2


def test_markitdown_placeholder_uses_docx_embedded_image_for_ocr(monkeypatch):
    from types import SimpleNamespace

    from PIL import Image
    from markitdown import MarkItDown

    from backend.documents.service import _markitdown_docx_text

    image_stream = io.BytesIO()
    Image.new("RGB", (420, 240), "white").save(image_stream, format="PNG")
    document = Document()
    document.add_heading("课程材料", level=1)
    document.add_picture(io.BytesIO(image_stream.getvalue()))
    stream = io.BytesIO()
    document.save(stream)

    monkeypatch.setattr(
        MarkItDown,
        "convert_stream",
        lambda self, *args, **kwargs: SimpleNamespace(
            text_content="# 课程材料\n\n![教学流程图](data:image/png;base64...)"
        ),
    )
    monkeypatch.setattr("backend.documents.service._cnocr_engine", lambda: object())

    def recognize(data: bytes, engine: object) -> list[str]:
        assert data.startswith(b"\x89PNG")
        return ["图中教学流程"]

    monkeypatch.setattr("backend.documents.service._ocr_image_bytes", recognize)

    text, _, image_count, warnings, ocr_engine_name = _markitdown_docx_text(stream.getvalue())

    assert "图片识别（教学流程图）" in text
    assert "图中教学流程" in text
    assert image_count == 1
    assert ocr_engine_name == "RapidOCR v6"
    assert any("识别 1 张" in warning for warning in warnings)


def test_parse_document_pdf_restores_reading_order_and_removes_repeated_margins():
    document = fitz.open()
    for index in range(3):
        page = document.new_page(width=595, height=842)
        page.insert_text((72, 35), "Thermodynamics Course", fontsize=9)
        page.insert_text((72, 100), "Energy Conservation" if index == 0 else f"Process Analysis {index}", fontsize=18)
        page.insert_text((72, 140), "Heat, work, and internal energy form a complete balance for the selected system.", fontsize=11)
        page.insert_text((290, 815), str(index + 1), fontsize=9)
    data = document.tobytes()
    document.close()

    result = parse_document("thermodynamics.pdf", data)

    assert result["extraction_report"]["engine"] == "PyMuPDF"
    assert result["extraction_report"]["page_count"] == 3
    assert result["raw_text"].count("Thermodynamics Course") <= 1
    assert "Energy Conservation" in [section["title"] for section in result["sections"]]
    assert result["sections"][0]["page_start"] == 1
    assert any("页眉" in warning for warning in result["extraction_report"]["warnings"])


def test_pdf_replaces_hidden_ocr_overlay_and_reports_every_page(monkeypatch):
    from PIL import Image

    from backend.documents.service import _has_hidden_ocr_overlay

    image_stream = io.BytesIO()
    Image.new("RGB", (595, 842), "white").save(image_stream, format="PNG")
    document = fitz.open()
    for page_number in range(2):
        page = document.new_page(width=595, height=842)
        page.insert_image(page.rect, stream=image_stream.getvalue())
        page.insert_text((72, 100), f"broken hidden OCR {page_number}", render_mode=3)
    assert _has_hidden_ocr_overlay(document[0]) is True
    data = document.tobytes()
    document.close()

    pages = [
        ("第2章 内置对象、运算符、表达式、关键字", "第一页准确正文"),
        ("2.1.1 常量与变量", "第二页准确正文"),
    ]

    def recognize(page, engine):
        title, body = pages[page.number]
        return [
            {"y": 80.0, "x": 60.0, "text": title, "size": 18.0, "bbox": (60, 80, 520, 102), "score": 0.99, "ocr": True, "line_index": 0},
            {"y": 130.0, "x": 60.0, "text": body, "size": 10.0, "bbox": (60, 130, 520, 144), "score": 0.97, "ocr": True, "line_index": 1},
        ]

    monkeypatch.setattr("backend.documents.service._preferred_ocr_engine", lambda: ("RapidOCR v6", object(), []))
    monkeypatch.setattr("backend.documents.service._ocr_page_elements", recognize)

    result = parse_document("hidden-ocr-book.pdf", data)

    assert "broken hidden OCR" not in result["raw_text"]
    assert "第一页准确正文" in result["raw_text"]
    assert "第二页准确正文" in result["raw_text"]
    assert result["extraction_report"]["ocr_page_count"] == 2
    assert [page["source_kind"] for page in result["extraction_report"]["page_reports"]] == ["hidden_ocr", "hidden_ocr"]
    assert [page["character_count"] > 0 for page in result["extraction_report"]["page_reports"]] == [True, True]
    section = next(item for item in result["sections"] if "常量与变量" in item["title"])
    assert section["page_start"] == 2
    assert any("隐藏 OCR" in warning for warning in result["extraction_report"]["warnings"])


def test_rapidocr_output_restores_visual_order_and_filters_low_confidence():
    from types import SimpleNamespace

    import numpy as np
    from PIL import Image

    from backend.documents.service import _ocr_image_bytes

    class RapidLikeEngine:
        def __call__(self, image):
            assert image.shape[:2] == (300, 400)
            return SimpleNamespace(
                txts=("右侧", "左侧", "低置信度噪声"),
                scores=(0.99, 0.98, 0.05),
                boxes=np.asarray([
                    [[220, 20], [300, 20], [300, 50], [220, 50]],
                    [[20, 20], [100, 20], [100, 50], [20, 50]],
                    [[20, 80], [160, 80], [160, 110], [20, 110]],
                ], dtype=np.float32),
            )

    stream = io.BytesIO()
    Image.new("RGB", (400, 300), "white").save(stream, format="PNG")

    assert _ocr_image_bytes(stream.getvalue(), RapidLikeEngine()) == ["左侧", "右侧"]


def test_preferred_ocr_engine_falls_back_to_cnocr(monkeypatch):
    from backend.documents.service import _preferred_ocr_engine

    fallback = object()

    def fail_rapidocr():
        raise ImportError("rapidocr unavailable")

    monkeypatch.setattr("backend.documents.service._rapidocr_engine", fail_rapidocr)
    monkeypatch.setattr("backend.documents.service._cnocr_engine", lambda: fallback)

    name, engine, warnings = _preferred_ocr_engine()

    assert name == "CnOCR"
    assert engine is fallback
    assert any("回退 CnOCR" in warning for warning in warnings)


def test_ocr_heading_rejects_decimal_values_but_keeps_numbered_titles():
    from backend.documents.service import _looks_like_ocr_heading

    assert _looks_like_ocr_heading("2.1.4列表、元组、字典、集合", 20) is True
    assert _looks_like_ocr_heading("2.3.6 map()函数、reduce()函数、filter()函数", 20) is True
    assert _looks_like_ocr_heading("3.14, 1.3e5", 20) is False
    assert _looks_like_ocr_heading("12.345", 20) is False


def test_fully_ocr_document_requires_review():
    from backend.documents.service import ExtractedDocument, _quality_report

    extracted = ExtractedDocument(
        text="课程正文" * 300,
        engine="PyMuPDF + CnOCR",
        page_count=10,
        title_count=10,
        image_count=10,
        text_block_count=10,
        scanned_page_count=10,
        ocr_page_count=10,
        source_line_count=100,
    )

    report = _quality_report(".pdf", extracted, extracted.text)

    assert report["quality_level"] == "medium"
    assert report["quality_score"] < 85


def test_knowledge_points_deduplicate_bulleted_numbered_heading():
    text = """# 第10章 深度学习\n## 动物视觉机理与深度学习的提出\n## 口 10.1 动物视觉机理与深度学习的提出\n"""

    titles = [point.title for point in extract_knowledge_points(text)]

    assert titles.count("动物视觉机理与深度学习的提出") == 1
